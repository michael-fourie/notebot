from pptx import Presentation, util, text
from pptx.util import Cm, Pt
from docx import Document


# Returns the pptx
def getPresentation(url):
    prs = Presentation(url)
    return prs


# Returns the text from the power point slides
def extractText(prs):

    text = []
    tempText = []
    slideCount = 0

    #Set up the array that will hold the text
    for i in prs.slides:
        text.append([])
        slideCount+=1

    slideCount = 0

    #Itterate through each slide, extract the text, and put it in our array
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    tempText.append(runElement)
        text[slideCount] = tempText
        tempText = []
        slideCount+=1

    return text


def makeDocument(text):
    document = Document()

    for slide in range(len(text)):
        for section in range(len(text[slide])):
            document.add_paragraph(text[slide][section])
    document.save('C:/Users/micha/PycharmProjects/notebot/backend/test.docx')


def main(url):
    prs = getPresentation(url)
    text = extractText(prs)
    makeDocument(text)


